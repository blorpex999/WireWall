from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "release" / "WireWall-Ydays-presentation.pptx"
LOGO = ROOT / "assets" / "wirewall_logo_128.png"

BG = RGBColor(15, 20, 27)
PANEL = RGBColor(22, 29, 39)
BORDER = RGBColor(48, 63, 80)
TEXT = RGBColor(237, 241, 246)
MUTED = RGBColor(153, 166, 184)
BLUE = RGBColor(40, 161, 255)
GREEN = RGBColor(103, 197, 135)
ORANGE = RGBColor(240, 190, 82)
RED = RGBColor(255, 101, 125)


SLIDES = [
    {
        "title": "WireWall",
        "subtitle": "Supervision USB locale pour poste Windows",
        "tag": "Projet Ydays",
        "bullets": ["Voir ce qui se connecte", "Evaluer le risque", "Garder une preuve exploitable"],
        "notes": "WireWall aide a transformer un branchement USB incertain en information lisible et exploitable.",
    },
    {
        "title": "Le probleme USB",
        "subtitle": "Une cle USB inconnue pose une question simple : que vient-on de brancher ?",
        "tag": "Probleme",
        "bullets": ["Visibilite limitee", "Risque difficile a qualifier", "Preuve difficile a retrouver"],
        "notes": "Le risque n'est pas seulement l'objet branche. Le probleme est le manque de visibilite et de tracabilite.",
    },
    {
        "title": "Notre reponse produit",
        "subtitle": "Observer, memoriser, evaluer, alerter, proposer",
        "tag": "Solution",
        "bullets": ["Application Windows locale", "Memoire SQLite du poste", "Decision humaine supervisee"],
        "notes": "WireWall ne promet pas une protection magique. Il rend la situation claire, suivable et defendable.",
    },
    {
        "title": "Comment ca marche",
        "subtitle": "Du branchement USB a l'incident",
        "tag": "Pipeline",
        "bullets": ["Detection USB -> Baseline -> Score", "Score -> Alerte -> Incident", "Suggestion -> Validation humaine"],
        "notes": "La baseline donne du contexte. Un peripherique nouveau, rare ou blacklist n'a pas le meme niveau de risque.",
    },
    {
        "title": "Ce que WireWall sait faire",
        "subtitle": "Les fonctionnalites utiles pour un poste Windows",
        "tag": "Produit",
        "bullets": ["Inventaire USB + score de risque", "Alertes, incidents, historique", "Controle USBSTOR + IA locale optionnelle"],
        "notes": "Le produit couvre le parcours complet : voir, qualifier, suivre et agir quand les droits Windows le permettent.",
    },
    {
        "title": "Reel, optionnel, limites",
        "subtitle": "Une promesse technique precise et defendable",
        "tag": "Honnetete",
        "bullets": ["Reel : monitoring, historique, alertes, incidents, USBSTOR", "Optionnel : IA locale via Ollama", "Limite : pas de driver noyau, USBSTOR bloque le stockage"],
        "notes": "Nous preferons une promesse precise plutot qu'un discours flou. WireWall indique ce qui est reel, deduit ou limite.",
    },
    {
        "title": "Architecture locale",
        "subtitle": "Tout reste sur le poste",
        "tag": "Technique",
        "bullets": ["PyQt6 pour l'interface", "PyUSB/libusb1 + SQLite + registre USBSTOR", "Ollama local seulement si l'IA est activee"],
        "notes": "L'architecture est volontairement locale : detection, memoire, exports et IA restent sur la machine.",
    },
    {
        "title": "Demo live",
        "subtitle": "Montrer la valeur, pas cliquer partout",
        "tag": "Demo",
        "bullets": ["Dashboard -> Peripheriques -> Alertes", "Controle USB -> Analyse IA", "Plan B si IA, admin ou USB reel manquent"],
        "notes": "On suit le chemin le plus utile : voir, comprendre, suivre, agir, puis terminer par l'IA si elle est disponible.",
    },
    {
        "title": "Impact et fiabilite",
        "subtitle": "Moins d'incertitude, plus de trace",
        "tag": "Preuves",
        "bullets": ["78 tests automatises valides", "Deux installateurs : standard et full", "Exports hashes et verification d'integrite"],
        "notes": "Le ROI est operationnel : gagner en visibilite, reduire l'incertitude et documenter les decisions.",
    },
    {
        "title": "Conclusion",
        "subtitle": "Un outil utile, local et defendable",
        "tag": "Synthese",
        "bullets": ["Visibilite USB", "Controle credible", "IA locale optionnelle et honnete"],
        "notes": "Notre objectif n'etait pas de tout promettre, mais de demontrer un outil utile, local et defendable.",
    },
    {
        "title": "Slide de secours",
        "subtitle": "Si une partie de la demo est degradee",
        "tag": "Backup",
        "bullets": ["Sans Ollama : le logiciel reste exploitable", "Sans admin : USBSTOR est limite, le monitoring reste visible", "Sans USB reel : le mode demo reste separe"],
        "notes": "A garder en fin de deck. Ne l'afficher que si la demo live est degradee.",
    },
]


def add_textbox(slide, left, top, width, height, text, size=24, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Segoe UI"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_panel(slide, left, top, width, height, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)
    return shape


def add_tag(slide, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.45), Inches(1.8), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text.upper()
    run.font.name = "Segoe UI"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_bullets(slide, bullets, accent):
    for index, bullet in enumerate(bullets):
        y = Inches(3.15 + index * 0.78)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.08), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = accent
        dot.line.color.rgb = accent
        add_textbox(slide, Inches(1.32), y, Inches(7.1), Inches(0.44), bullet, size=21, color=TEXT)


def add_notes(slide, notes):
    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    accents = [BLUE, ORANGE, BLUE, GREEN, BLUE, RED, BLUE, ORANGE, GREEN, BLUE, RED]

    for idx, data in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG

        accent = accents[idx]
        add_tag(slide, data["tag"], accent)
        add_textbox(slide, Inches(0.72), Inches(0.48), Inches(6.8), Inches(0.38), f"{idx + 1:02d}", 12, MUTED, True)
        add_textbox(slide, Inches(0.72), Inches(0.9), Inches(8.4), Inches(0.72), data["title"], 34, TEXT, True)
        add_textbox(slide, Inches(0.75), Inches(1.65), Inches(8.7), Inches(0.58), data["subtitle"], 18, MUTED)

        add_panel(slide, Inches(0.78), Inches(2.75), Inches(8.7), Inches(3.4))
        add_bullets(slide, data["bullets"], accent)

        add_panel(slide, Inches(9.85), Inches(1.28), Inches(2.78), Inches(4.85))
        if LOGO.exists():
            slide.shapes.add_picture(str(LOGO), Inches(10.78), Inches(1.72), Inches(0.95), Inches(0.95))
        add_textbox(slide, Inches(10.12), Inches(2.95), Inches(2.25), Inches(0.55), "WireWall", 22, TEXT, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(10.08), Inches(3.55), Inches(2.32), Inches(1.0), "Voir\nComprendre\nProuver", 16, MUTED, False, PP_ALIGN.CENTER)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.color.rgb = accent
        add_notes(slide, data["notes"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
